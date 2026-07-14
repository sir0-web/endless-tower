# -*- coding: utf-8 -*-
"""KIZUNA concept deck generator — EbT-styled (dark navy / gold / purple)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (from EbT_Newgeneration.png) ----
BG      = RGBColor(0x0D, 0x0A, 0x1A)   # near-black navy
PANEL   = RGBColor(0x1C, 0x15, 0x33)   # dark purple-navy panel
PANEL2  = RGBColor(0x14, 0x10, 0x28)
GOLD    = RGBColor(0xF0, 0xC6, 0x4A)
GOLD_D  = RGBColor(0xC9, 0x9A, 0x2C)
PURPLE  = RGBColor(0x9B, 0x5C, 0xF6)
RED     = RGBColor(0xC8, 0x33, 0x3C)
TEXT    = RGBColor(0xED, 0xE6, 0xD8)   # warm off-white
MUTED   = RGBColor(0xB9, 0xAF, 0xD0)
GREEN   = RGBColor(0x7E, 0xD3, 0x8A)
GRAY    = RGBColor(0x8A, 0x84, 0x9C)

FONT_TITLE = "Yu Mincho Demibold"
FONT_BODY  = "Yu Gothic UI"

SW, SH = Inches(13.333), Inches(7.5)


def set_ea_font(run, name):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # thin purple glow line at very top
    g = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Pt(3))
    g.fill.solid(); g.fill.fore_color.rgb = PURPLE
    g.line.fill.background(); g.shadow.inherit = False
    return s


def txt(s, l, t, w, h, runs, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
        font=FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs: str, or list of paragraphs; each paragraph is str or list of (text, dict)."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        if isinstance(para, str):
            para = [(para, {})]
        for text, ov in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
            f.name = ov.get("font", font)
            set_ea_font(r, ov.get("font", font))
    return tb


def panel(s, l, t, w, h, fill=PANEL, line=GOLD_D, line_w=1.2, round_=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, l, t, w, h)
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def header(s, title, sub=None):
    txt(s, Inches(0.55), Inches(0.28), Inches(11.0), Inches(0.75),
        title, size=30, color=GOLD, bold=True, font=FONT_TITLE)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.02),
                            Inches(12.2), Pt(1.6))
    ln.fill.solid(); ln.fill.fore_color.rgb = GOLD_D
    ln.line.fill.background(); ln.shadow.inherit = False
    if sub:
        txt(s, Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.4),
            sub, size=13, color=MUTED)


def chip(s, l, t, w, text, color=GOLD):
    c = panel(s, l, t, w, Inches(0.42), fill=PANEL2, line=color, line_w=1.0)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color
    r.font.name = FONT_BODY
    set_ea_font(r, FONT_BODY)
    c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c


# ============ 1. Title ============
s = slide()
# purple band behind title
band = panel(s, 0, Inches(2.15), SW, Inches(2.7), fill=PANEL2, line=None, round_=False)
txt(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.5),
    "Endless Basement Tower — 新章企画", size=18, color=PURPLE, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(2.35), Inches(11.7), Inches(1.5),
    "KIZUNA", size=88, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.6),
    "一人では攻略できない塔を、仲間の“意志”を受け継ぎながら全員で踏破する",
    size=18, color=TEXT, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.6),
    "非同期協力型ローグライク", size=24, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
    "企画コンセプト・仕様まとめ｜出典：まっち×モコ 打ち合わせ（2026-07-07〜07-12）",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ============ 2. コンセプト（北極星） ============
s = slide()
header(s, "1. コンセプト（確定・北極星）", "迷ったらここに立ち返って方向修正する")
p = panel(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.35), line=GOLD, line_w=2)
txt(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(1.0),
    "「一人では攻略できない塔を、仲間の“意志”を受け継ぎながら\n全員で踏破する非同期協力型ローグライク」",
    size=21, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(0.75), Inches(3.15), Inches(6.0), Inches(0.4), "■ 狙い",
    size=16, color=PURPLE, bold=True)
txt(s, Inches(0.95), Inches(3.6), Inches(6.6), Inches(2.6), [
    "・プレイヤーはそれぞれ自分のタイミングで塔へ挑戦する",
    [("・", {}), ("敗北は失敗ではなく、次の挑戦者へ力を託す行為", {"color": GOLD, "bold": True}), ("となる", {})],
    "・誰かの犠牲が、誰かの前進になり、最終的にはギルド全員の力で塔を攻略する",
], size=15)
txt(s, Inches(8.0), Inches(3.15), Inches(4.0), Inches(0.4), "■ キーワード",
    size=16, color=PURPLE, bold=True)
kws = ["非同期協力", "死亡＝貢献", "継承", "ローグライク", "コミュニティ攻略", "リレー形式"]
for i, kw in enumerate(kws):
    chip(s, Inches(8.0 + (i % 2) * 2.3), Inches(3.6 + (i // 2) * 0.6), Inches(2.1), kw)
txt(s, Inches(8.0), Inches(5.5), Inches(4.4), Inches(0.5),
    "＝「次へ繋ぐ」ゲーム体験", size=17, color=GOLD, bold=True)

# ============ 3. プレイヤー体験の転換 ============
s = slide()
header(s, "2. プレイヤー体験の転換")
panel(s, Inches(0.7), Inches(1.7), Inches(5.6), Inches(4.6), fill=PANEL2, line=GRAY)
txt(s, Inches(0.7), Inches(1.95), Inches(5.6), Inches(0.5), "従来のゲーム",
    size=19, color=GRAY, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(1.1), Inches(2.7), Inches(4.8), Inches(2.4), [
    "戦う", "　↓", "強くなる", "　↓", "自分でクリアする",
], size=16, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(5.45), Inches(5.2), Inches(0.6),
    "自分が強くなってクリアする。", size=15, color=MUTED, align=PP_ALIGN.CENTER, bold=True)

panel(s, Inches(7.0), Inches(1.7), Inches(5.6), Inches(4.6), line=GOLD, line_w=2)
txt(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.5), "KIZUNA（新章）",
    size=19, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(7.3), Inches(2.7), Inches(5.0), Inches(2.5), [
    "力尽きても次の仲間へ力を託す", "　↓", "受け取った力でさらに前へ", "　↓",
    "みんなの意志を繋いで頂上へ", "　↓",
    [("ギルド全員の力で塔を攻略", {"color": GOLD, "bold": True})],
], size=15, align=PP_ALIGN.CENTER, line_spacing=1.0)
txt(s, Inches(6.9), Inches(5.6), Inches(5.8), Inches(0.6),
    "自分が次の挑戦者を強くして、みんなでクリアする。",
    size=15, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
# arrow between panels
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.35), Inches(3.7), Inches(0.6), Inches(0.6))
ar.fill.solid(); ar.fill.fore_color.rgb = PURPLE; ar.line.fill.background(); ar.shadow.inherit = False

# ============ 4. 体験イメージ ============
s = slide()
header(s, "3. 体験イメージ — 意志のリレー", "コンセプトビジュアルより")
msgs = [
    ("「15Fで力尽きた…でも、宝と経験を残しておくよ！」", "− 10分前"),
    ("「ありがとう！その力で16Fを突破したよ！」", "− 30分前"),
    ("「おかげで20Fまで来れた！次はきっと頂上へ…！」", "− 1時間前"),
]
for i, (m, when) in enumerate(msgs):
    t = Inches(1.7 + i * 1.35)
    panel(s, Inches(1.2 + i * 0.7), t, Inches(9.2), Inches(1.0), line=GOLD_D)
    txt(s, Inches(1.5 + i * 0.7), t + Inches(0.12), Inches(8.6), Inches(0.5),
        m, size=17, color=TEXT, bold=True)
    txt(s, Inches(1.5 + i * 0.7), t + Inches(0.55), Inches(8.6), Inches(0.35),
        when, size=12, color=MUTED)
    if i < 2:
        a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                               Inches(5.6 + i * 0.7), t + Inches(1.02), Inches(0.4), Inches(0.3))
        a.fill.solid(); a.fill.fore_color.rgb = PURPLE
        a.line.fill.background(); a.shadow.inherit = False
txt(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.7),
    "—— みんなの意志が、塔を登っていく。", size=24, color=GOLD, bold=True,
    align=PP_ALIGN.CENTER, font=FONT_TITLE)

# ============ 5. 新章の位置づけ ============
s = slide()
header(s, "4. 新章の位置づけ（確定）")
rows = [
    ("EbT内の1コンテンツ", "別ゲームではなくEbT内に制作。TOPのモード選択 or ひろばから「オンラインの世界へ」（仮）を選んで新章へ転生する入口イメージ"),
    ("やり切り型（クリアあり）", "協力クリア型ローグライク。エンドレス型の本編とはここが根本的に異なる。アップデートで上層階を追加する余地あり"),
    ("ストーリーモードあり", "コンテンツ名称で本編と切り分け。理想の体験アーク＝ストーリー性 → クリアへの盛り上がり → エンドロールの達成感"),
    ("マップは現行仕様を流用", "マップ・基礎構成は現行EbTの仕様にまるまる乗せる（新規マップシステムの開発はしない）"),
]
for i, (h, b) in enumerate(rows):
    t = Inches(1.6 + i * 1.35)
    panel(s, Inches(0.7), t, Inches(3.3), Inches(1.15), line=GOLD_D)
    txt(s, Inches(0.85), t + Inches(0.1), Inches(3.0), Inches(0.95),
        h, size=15, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.3), t + Inches(0.08), Inches(8.4), Inches(1.1),
        b, size=13.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ============ 6. 継承システム 4類型 ============
s = slide()
header(s, "5. 継承システム — 4類型の採否", "方向性確定：③＋④の組み合わせが核")
rows = [
    ("① ステータス合算型", "死んだ仲間のステータスの一部を次の仲間へ永久加算", "△ 主軸にしない",
     "「他者への貢献」を表現しにくい。ギルドツリー解禁後の微増ボーナス扱い", GRAY),
    ("② スキル選択・融合型", "死者の必殺技/パッシブを遺品・魂として選んでセット", "△ 初期は見送り",
     "物語性は最強だがコスト重。将来「形見のパッシブ1つ」として③④へ統合", GRAY),
    ("③ ツリー解禁型（家系図）", "死者の実績ポイントで一族全体の強化ツリーをアンロック", "◎ 採用",
     "一族＝ギルド全体と読み替え。誰の死も共有ツリーに貢献＝「死亡＝貢献」が自然に成立", GOLD),
    ("④ リソース限定型（形見）", "死者の装備1つ or 最高ステータス1つを選択して継承", "◎ 採用",
     "ドッペルゲンガーと直結。「誰かの犠牲→誰かの前進」を1アクションで体験", GOLD),
]
for i, (name, desc, verdict, why, col) in enumerate(rows):
    t = Inches(1.65 + i * 1.32)
    panel(s, Inches(0.7), t, Inches(11.9), Inches(1.15),
          line=col, line_w=1.8 if col == GOLD else 1.0)
    txt(s, Inches(0.9), t + Inches(0.08), Inches(2.9), Inches(1.0),
        name, size=14, color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.85), t + Inches(0.08), Inches(3.6), Inches(1.0),
        desc, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.55), t + Inches(0.08), Inches(1.45), Inches(1.0),
        verdict, size=13, color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(9.05), t + Inches(0.08), Inches(3.4), Inches(1.0),
        why, size=11, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

# ============ 7. ③ギルドツリー + ④形見セレクト ============
s = slide()
header(s, "6. 採用した2つの継承の中身")
panel(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(5.3), line=GOLD, line_w=1.8)
txt(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(0.5),
    "③ ギルド共有「塔攻略ツリー」", size=18, color=GOLD, bold=True,
    align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(1.05), Inches(2.6), Inches(5.3), Inches(4.0), [
    [("・誰が死んでも", {}), ("貢献値がギルド全体に加算 → ツリー解禁", {"bold": True, "color": GOLD})],
    "・「みんなで攻略する」の骨格",
    "・ツリー内容の候補：固定の全員共有ステータス or パッシブスキル（モコ案）",
    "・転生後のフルポイント分岐で個人のプレイスタイルの幅も担保（要詳細化）",
], size=14)
panel(s, Inches(6.95), Inches(1.6), Inches(5.9), Inches(5.3), line=GOLD, line_w=1.8)
txt(s, Inches(6.95), Inches(1.85), Inches(5.9), Inches(0.5),
    "④ 形見セレクト（ドッペル拡張）", size=18, color=GOLD, bold=True,
    align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(7.3), Inches(2.6), Inches(5.3), Inches(4.2), [
    "・既存ドッペルゲンガー（死亡キャラ登録 → 他PLの±10階に出現 → 撃破）の上に載せる",
    "・撃破時にモーダルで1つ選択する形へ拡張：",
    [("　1. 装備1つ", {"color": GOLD}), ("（equipmentから選択）", {"color": MUTED, "size": 12})],
    [("　2. 最高ステータス1点", {"color": GOLD}), ("（str/agi/dex/int/vit/lukの最大値）", {"color": MUTED, "size": 12})],
    [("　3. stat_point_reward の一部 or 全部", {"color": GOLD})],
    "・equipmentはDB保存済みのため追加コスト小",
    [("・モコ：「まさに理想」（確定）", {"color": GREEN, "bold": True})],
], size=14)

# ============ 8. 魂のメッセージ ============
s = slide()
header(s, "7. 魂のメッセージ（確定＋詳細方針）")
panel(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(2.5), line=GOLD_D)
txt(s, Inches(0.95), Inches(1.75), Inches(5.4), Inches(0.4), "■ 基本仕様",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(1.0), Inches(2.2), Inches(5.4), Inches(1.9), [
    [("・死亡時、その場（死亡地点）にメッセージを残せる", {"bold": True})],
    "・用途：応援／「この先のXXXやべーぞ！」的な警告・攻略情報",
    "・画面イメージはデバッグ環境で作成済み（07-12）",
], size=13)
panel(s, Inches(0.7), Inches(4.35), Inches(5.9), Inches(2.5), line=GOLD_D)
txt(s, Inches(0.95), Inches(4.5), Inches(5.4), Inches(0.4), "■ 増殖・容量対策（方針決定）",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(1.0), Inches(4.95), Inches(5.4), Inches(1.85), [
    "・文字数とエフェクト整理後、新規50体ほど保存できる規模に調整",
    [("・削除は時限式（48h or 72h）", {"bold": True}), ("が基本方針", {})],
    "・「触れた回数で削除」は現接続者数では残り続けるため不採用",
], size=13)
panel(s, Inches(6.95), Inches(1.6), Inches(5.9), Inches(5.25), line=PURPLE)
txt(s, Inches(7.2), Inches(1.75), Inches(5.4), Inches(0.4), "■ 発展アイデア（案として保持・要設計）",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(7.25), Inches(2.25), Inches(5.35), Inches(4.5), [
    [("魂のログ（記録書）：", {"color": GOLD, "bold": True}),
     ("触れられたメッセージが全員共有の書物に記録され、いつでも閲覧できる", {})],
    [("　⚠ 「見つける意義」「シェアするメリット」の設計必須", {"color": RED, "size": 12})],
    [("ゲーム外での閲覧：", {"color": GOLD, "bold": True}),
     ("セットアップ画面から死亡時ログを振り返る「思い出」機能（モコ案）", {})],
    [("クリア後エンドロール：", {"color": GOLD, "bold": True}),
     ("死亡ログとメッセージがスタッフロールのように流れる。実装容易（まっち確認済み）。「エンドロールの達成感」と直結する有力候補", {})],
], size=13)

# ============ 9. 拠点成長システム ============
s = slide()
header(s, "8. 拠点成長システム（箱庭）— 追加要素として評価済み",
       "スタート地点の「ふもとの村」を共有スペースとして協力育成（モコ提案 07-12）")
txt(s, Inches(0.75), Inches(1.6), Inches(5.9), Inches(0.4), "■ 施設候補（仮）",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(0.95), Inches(2.05), Inches(5.6), Inches(1.0), [
    "強化施設：武器屋・防具屋・鍛冶屋・雑貨屋・精錬など",
    "資源施設：畑・鉱山・森など",
], size=13)
txt(s, Inches(0.75), Inches(3.1), Inches(5.9), Inches(0.4), "■ 解放条件の評価（仮）",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(0.95), Inches(3.55), Inches(5.6), Inches(2.2), [
    [("◎ 進行度による解放", {"color": GOLD}), ("　— 王道", {"color": MUTED, "size": 12})],
    [("◎ ダンジョン内で発見して解放（最推奨）", {"color": GOLD, "bold": True})],
    [("　　「捕まっている大工を助けて村へ送る」等。探索動機と直結", {"color": MUTED, "size": 12})],
    [("◎ 資源投入で解放", {"color": GOLD}), ("　— 発見型とセットで機能", {"color": MUTED, "size": 12})],
    [("✕ プレイ回数による解放", {"color": RED}), ("　— 自殺周回で破綻", {"color": MUTED, "size": 12})],
], size=13)
txt(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(0.4), "■ 成長システムの評価（仮）",
    size=15, color=PURPLE, bold=True)
txt(s, Inches(7.2), Inches(2.05), Inches(5.5), Inches(2.3), [
    [("◎ 死亡時に進行度ポイント付与 → 投入で成長（最推奨）", {"color": GOLD, "bold": True})],
    [("　　深いほど貢献大＝「死亡＝貢献」と完全一致。③ギルドツリーの貢献値と統合できる可能性", {"color": MUTED, "size": 12})],
    [("○ 資源を一定投入で成長", {}), ("　— 資源システムの新規実装が前提", {"color": MUTED, "size": 12})],
    [("△ 時間による成長", {}), ("　— オフライン中の緩回復程度なら可", {"color": MUTED, "size": 12})],
    [("✕ プレイ回数", {"color": RED}), ("　— 同上", {"color": MUTED, "size": 12})],
], size=13)
panel(s, Inches(7.0), Inches(4.75), Inches(5.7), Inches(2.05), line=GOLD_D)
txt(s, Inches(7.2), Inches(4.9), Inches(5.3), Inches(0.4), "■ 導入方針（まっち所感）",
    size=14, color=GOLD, bold=True)
txt(s, Inches(7.25), Inches(5.3), Inches(5.3), Inches(1.5), [
    [("フェーズ1：「捕まっている大工を助ける」型の個人の町拡張をまず導入して様子見", {"bold": True})],
    "協力・共闘の深化と現状の単独プレイ傾向のギャップを埋めるプロモーションも併せて必要",
], size=12.5)

# ============ 10. 設計上の要注意ポイント ============
s = slide()
header(s, "9. 設計上の要注意ポイント（実装前に解決）")
rows = [
    ("ランキングの公平性（最重要）", RED,
     "継承・拠点の恩恵でスタート時の強さが変わると歴代記録と比較不能に。対策：a) 恩恵を利便性に限定 b) 村Lv併記/シーズン制。新章は別コンテンツなので新章専用の進行度表示にすれば本編ランキングとの衝突は回避可能"),
    ("クライアント信頼問題", GOLD,
     "ゲームロジックは全てクライアント側でチート耐性なし。共有進捗への貢献はサーバー側RPCで上限・レート制御が必須（ジャックポット冪等化と同型の設計で対応可能）"),
    ("自殺周回対策", GOLD,
     "貢献量は「回数」ではなく「深度」を通貨にする（浅い死は貢献が小さい設計）"),
    ("イベント競合", GOLD,
     "すかるぽりん等の既存全鯖イベントとの同時発生時の挙動整理（過去に湧き時間20分固定の制約で移動ロジックを断念した経緯あり）"),
]
for i, (h, col, b) in enumerate(rows):
    t = Inches(1.6 + i * 1.35)
    panel(s, Inches(0.7), t, Inches(11.9), Inches(1.18), line=col,
          line_w=2 if col == RED else 1.2)
    txt(s, Inches(0.9), t + Inches(0.08), Inches(3.1), Inches(1.0),
        h, size=15, color=col, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.15), t + Inches(0.06), Inches(8.3), Inches(1.06),
        b, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ============ 11. 既存システムとの関係 ============
s = slide()
header(s, "10. 既存システムとの関係 — KIZUNAの土台", "新規開発を最小化できる根拠")
rows = [
    ("ドッペルゲンガー", "④形見セレクトの土台。死亡キャラ登録・出現ロジックをほぼ流用"),
    ("ワールド通知・いいね", "継承の発生・お礼の伝達に流用（「○○さんの遺産で△△さんが16F突破！」→いいね）"),
    ("ジャックポット（冪等化済みRPC）", "共有ゲージ・貢献受付のサーバー側実装の雛形"),
    ("クラウドセーブ", "プレイヤー識別の既存基盤（名前＋パスワード）"),
    ("掲示板（ひろば）", "コミュニティ攻略の場。新章への入口候補でもある"),
    ("現行マップ・ダンジョン生成", "新章のマップはこれをそのまま使用（新規開発なし）"),
]
for i, (h, b) in enumerate(rows):
    t = Inches(1.65 + i * 0.9)
    panel(s, Inches(0.7), t, Inches(3.9), Inches(0.75), line=GOLD_D)
    txt(s, Inches(0.88), t + Inches(0.04), Inches(3.6), Inches(0.68),
        h, size=13.5, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.9), t + Inches(0.04), Inches(7.8), Inches(0.7),
        b, size=13, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ============ 12. 進め方・体制 ============
s = slide()
header(s, "11. 進め方・体制（確定）")
txt(s, Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.4), "■ 順番が大切",
    size=16, color=PURPLE, bold=True)
txt(s, Inches(0.95), Inches(2.05), Inches(11.5), Inches(1.1), [
    [("根幹となるメインコンテンツをフォーカスして優先実装 → 関連コンテンツを芋づる式で順に実装", {"bold": True})],
    "そのための相関図を作成中（Claude整理）。「イイ感じ」になったら関係性の強いところから着手",
], size=14)
txt(s, Inches(0.75), Inches(3.3), Inches(11.8), Inches(0.4), "■ スケジュール感（2026-07-12時点）",
    size=16, color=PURPLE, bold=True)
steps = ["弓実装（完了）", "細かい修正", "7月中旬〜 新機能着手"]
for i, st in enumerate(steps):
    l = Inches(1.0 + i * 3.9)
    c = panel(s, l, Inches(3.8), Inches(3.2), Inches(0.7),
              line=GOLD if i == 2 else GOLD_D, line_w=1.8 if i == 2 else 1.0)
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
    r = pgh.add_run(); r.text = st
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = GOLD if i == 2 else TEXT
    r.font.name = FONT_BODY; set_ea_font(r, FONT_BODY)
    if i < 2:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l + Inches(3.3), Inches(3.95),
                               Inches(0.5), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = PURPLE
        a.line.fill.background(); a.shadow.inherit = False
txt(s, Inches(0.75), Inches(4.95), Inches(11.8), Inches(0.4), "■ 役割分担・作画",
    size=16, color=PURPLE, bold=True)
txt(s, Inches(0.95), Inches(5.4), Inches(11.5), Inches(1.6), [
    [("まっち：", {"color": GOLD, "bold": True}), ("実装・設計・実現性判断・企画書への集約", {})],
    [("モコ：", {"color": GOLD, "bold": True}),
     ("現実性を考慮しないアイデア・ユーザー視点の理想形（システム/ストーリー/職業/アイテム/グラフィック）。ツリースキルのアイデア等を思考中", {})],
    [("作画：", {"color": GOLD, "bold": True}), ("base数点＋生成AIで負担を抑えて実装", {})],
], size=14)

# ============ 13. 未決事項 ============
s = slide()
header(s, "12. 未決事項（次回以降の議論テーマ）")
items = [
    "ギルドツリーの中身（解禁ノード具体案・フルポイント分岐）",
    "貢献値の算出式（深度ベース。到達階数×係数？）",
    "「ギルド」の単位（全鯖1共同体 or 複数ギルド制）",
    "「塔の攻略（クリア）」の定義（目標階数・章構成・周回）",
    "ストーリー原案・チュートリアル（死亡と引き継ぎを最初に体感）",
    "魂メッセージのモデレーション（不適切ワード対策）",
    "魂のログの「見つける意義・シェアするメリット」設計",
    "エンドロール演出の詳細（死亡ログの範囲・順序）",
    "新章と本編の進行データの関係（転生時の持ち込み）",
    "同時プレイ時（オンライン重複時）の楽しみの工夫",
    "協力プレイへ誘導するプロモーション設計",
]
for i, it in enumerate(items):
    col_i, row_i = i % 2, i // 2
    l = Inches(0.7 + col_i * 6.15)
    t = Inches(1.65 + row_i * 0.88)
    panel(s, l, t, Inches(5.95), Inches(0.72), fill=PANEL2, line=GOLD_D, line_w=0.75)
    txt(s, l + Inches(0.2), t + Inches(0.03), Inches(5.6), Inches(0.66),
        [[("□ ", {"color": GOLD, "bold": True}), (it, {})]],
        size=12.5, anchor=MSO_ANCHOR.MIDDLE)

# ============ 14. Closing ============
s = slide()
txt(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2),
    "みんなの意志が、塔を登っていく。", size=40, color=GOLD, bold=True,
    align=PP_ALIGN.CENTER, font=FONT_TITLE)
txt(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.6),
    "KIZUNA — 「次へ繋ぐ」ゲーム体験", size=18, color=PURPLE, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
    "最終更新：2026-07-12｜出典：Discordログ全文とコンセプト画像", size=11,
    color=MUTED, align=PP_ALIGN.CENTER)

out = r"C:\Users\USER\Desktop\endless-tower\docs\KIZUNA_concept_slides.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
